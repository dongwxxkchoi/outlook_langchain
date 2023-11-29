import win32com.client
import os
import glob
import olefile

from urllib.parse import unquote
import email.charset
from email.message import EmailMessage

import zipfile
import docx
from pptx import Presentation
import pandas as pd
import chromadb
from chromadb.utils import embedding_functions

from langchain.schema.document import Document
from langchain.document_loaders import TextLoader, PyPDFLoader, CSVLoader
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores.chroma import Chroma


import email
import re 
import poplib
import chardet
import base64
from datetime import datetime
from bs4 import BeautifulSoup

import sys
sys.path.append('C:/Users/dwchoi0610/dongwook/langchain/outlook/outlook_pipeline/src')

os.environ['openai_api_key'] = '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'

attachment_folder = 'C:/Temp/outlook/attachments'
eml_folder = 'C:/Temp/outlook/eml'

def delete_files(folder:str, extension: str):
    for f in glob.glob(f"{folder}/*.{extension}"):
        os.remove(f)


def fetch_outlook():
    # fetch application from outlook application
    outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
    # 6 means INBOX folder
    inbox = outlook.GetDefaultFolder(6)
    return inbox.Items


def pop_outlook():
    EMAIL_ADDRESS = "dwchoi0610@outlook.com"
    EMAIL_PASSWORD = "Outlookdongwook1!"

    Mailbox = poplib.POP3_SSL("outlook.office365.com", 995)
    Mailbox.user(EMAIL_ADDRESS)
    Mailbox.pass_(EMAIL_PASSWORD)

    return Mailbox


def decode_header(message):
    extracted, charset = email.header.decode_header(message)[0]

    # when encoding is specified
    if charset:
        extracted = extracted.decode(charset)
        return extracted
    else:
        # but if extracted is still bytes
        # maybe ascii and utf-8 encoded text is mixed
        if isinstance(extracted, bytes):
            tokens = message.split(' ')
            text = ""
            for token in tokens:
                token_extracted, token_charset = email.header.decode_header(token)[0]
                if token_charset:
                    text += token_extracted.decode(token_charset) + " "
                else:
                    text += token + " "
            return text
        else:
            return extracted


def parse_email_body(msg):
    content_type = msg.get_content_type().lower()
    if content_type=='text/plain' or content_type=='text/html':
        # get text content.
        content = msg.get_payload(decode=True)
        # get text charset.
        charset = msg.get_charset()
        # if can not get charset. 

        if charset is None:
            # get message 'Content-Type' header value.
            content_type = msg.get('Content-Type', '').lower()
            # parse the charset value from 'Content-Type' header value.
            pos = content_type.find('charset=')
            if pos >= 0:
                charset = content_type[pos + 8:].strip()
                pos = charset.find(';')
                if pos>=0:
                    charset = charset[0:pos]           
        if charset:
            content = content.decode(charset)

        if 'html' in content_type:
            soup = BeautifulSoup(content.strip(), 'html.parser')
            text = soup.get_text(separator=' ', strip=True)
            return text
        else:
            return content

    elif content_type.startswith('multipart'):
        # get multiple part list.
        body_msg_list = msg.get_payload()
        # loop in the multiple part list.
        for body_msg in body_msg_list:
            # parse each message part.
            return parse_email_body(body_msg)


def parse_email_content(attachment_paths, msg):
    # get message content type.
    content_type = msg.get_content_type().lower()
    
    # print('---------------------------------' + content_type + '------------------------------------------')
    # if the message part is text part.
    if content_type=='text/plain' or content_type=='text/html':
        pass

    # if this message part is still multipart such as 'multipart/mixed','multipart/alternative','multipart/related'
    elif content_type.startswith('multipart'):
        # get multiple part list.
        body_msg_list = msg.get_payload()
        # loop in the multiple part list.
        for body_msg in body_msg_list:
            # parse each message part.
            parse_email_content(attachment_paths, body_msg)

    # if this message part is an attachment part that means it is a attached file.        
    elif content_type.startswith('image') or content_type.startswith('application'):
        if content_type.startswith('image'):
            # not parse image
            return

        # get message header 'Content-Disposition''s value and parse out attached file name.
        attach_file_info_string = msg.get('Content-Disposition')
        prefix = '?utf-8?B?'

        if prefix in attach_file_info_string:
            matches = re.findall(r'=\?utf-8\?B\?(.*?)\?=', attach_file_info_string)
            decoded_text = ""

            for match in matches:
                decoded_bytes = base64.b64decode(match)
                decoded_text += decoded_bytes.decode('utf-8')
        
            attach_file_name = decoded_text

        else:
            start_pos = attach_file_info_string.find("filename=\"")
            end_pos = attach_file_info_string.find("\";")
            attach_file_name = attach_file_info_string[start_pos + len("filename=\""): end_pos]
        
        # get attached file content.
        attach_file_data = msg.get_payload(decode=True)

        # get current script execution directory path. 
        current_path = os.path.dirname(os.path.abspath(__file__))

        # get the attached file full path.
        attach_file_path = current_path + '/' + attach_file_name

        # write attached file content to the file.
        with open(attach_file_path,'wb') as f:
            f.write(attach_file_data)
        
        print('attached file is saved in path ' + attach_file_path)   
        attachment_paths.append(attach_file_path)
                
    else:
        content = msg.as_string()
        print(content)


def export_data_pop():
    Mailbox = pop_outlook()
    numMessages = len(Mailbox.list()[1])
    
    # langchain docs
    docs = list()
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    for i in range(numMessages):
        raw_email = b'\n'.join(Mailbox.retr(i+1)[1])
        msg = email.message_from_bytes(raw_email)
        
        # find header 
        msg_subject = decode_header(msg['Subject'])
        msg_from = decode_header(msg['From'])
        msg_to = decode_header(msg['To'])
        msg_date = decode_header(msg['Date'])

        metadata = {
            "Subject": str(msg_subject),
            "From": str(msg_from),
            "To": str(msg_to),
            "Date": str(datetime.strptime(msg_date, "%a, %d %b %Y %H:%M:%S %z")),
            "Attchment": False
        }

        # find body
        ## find text/plain data
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                cdispo = str(part.get('Content-Disposition'))
                # text/plain data without attachment 
                if ctype == 'text/plain' and 'attachment' not in cdispo:
                    body = part.get_payload(decode=True) 
                    charset = part.get_content_charset()
                    break

                # 아니라면, msg 전체적으로 payload 호출해서 넘겨줌
                body = msg.get_payload(decode=True)
                charset = part.get_content_charset()
        else:
            body = msg.get_payload(decode=True)
            charset = part.get_content_charset()

        ## get text/plain data. if not find text/html data
        try:
            msg_body = body.decode(charset)
        except Exception:
            msg_body = parse_email_body(msg)

        docs.extend([Document(page_content=x, metadata = metadata) for x in text_splitter.split_text(body)])

        # find attachment
        attachment_paths= []
        msg_attachments = parse_email_content(attachment_paths, msg)


def make_folders():
    try:
        os.makedirs('C:/Temp/outlook/attachments')
    except Exception:
        print("attachment folder already created")


def attachment_parsing(docs, attachment_path, filename, metadata, text_splitter):

    if filename.endswith(".pdf"):
        loader = PyPDFLoader(attachment_path)
        loaded = loader.load()
        for i, doc in enumerate(loaded):
            loaded[i].metadata.update(metadata)

        if isinstance(loaded[0], list):
            print(filename)
        
        docs.extend(loaded)

    elif filename.endswith((".docx", ".doc")):
        doc = docx.Document(attachment_path)
        text = ' '.join([paragraph.text for paragraph in doc.paragraphs])
        
        if isinstance([Document(page_content=x, metadata = metadata) for x in text_splitter.split_text(text)][0], list):
            print(filename)
        
        docs.extend([Document(page_content=x, metadata = metadata) for x in text_splitter.split_text(text)])

    elif filename.endswith((".xlsx", ".xls")):
        df = pd.read_excel(attachment_path)
        filename = filename.split('.')[0]+".txt"
        txt_path = attachment_path.split('.')[0] + ".txt"
        df.to_csv(txt_path, sep='\t', index=False, encoding='cp949')
        attachment_parsing(docs, txt_path, filename, metadata, text_splitter)

    elif filename.endswith(".csv"):
        loader = CSVLoader(attachment_path)
        loaded = loader.load()
        for i, doc in enumerate(loaded):
            loaded[i].metadata.update(metadata)

        if isinstance(loaded[0], list):
            print(filename)
        
        docs.extend(loaded)

    elif filename.endswith(".hwp"):
        f = olefile.OleFileIO(attachment_path)  
        encoded_text = f.openstream('PrvText').read() 
        decoded_text = encoded_text.decode('utf-16')

        if isinstance([Document(page_content=x, metadata = metadata) for x in text_splitter.split_text(decoded_text)][0], list):
            print(filename)
        
        docs.extend([Document(page_content=x, metadata=metadata) for x in text_splitter.split_text(decoded_text)])

    elif filename.endswith((".pptx", ".ppt")):
        prs = Presentation(attachment_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        if isinstance([Document(page_content=x, metadata=metadata) for x in text_splitter.split_text(text)][0], list):
            print(filename)
        
        docs.extend([Document(page_content=x, metadata=metadata) for x in text_splitter.split_text(text)])

    elif filename.endswith(".zip"):
        with zipfile.ZipFile(attachment_path, 'r') as zip_ref:
            file_names = zip_ref.namelist()
            zip_ref.extractall(attachment_folder)
            for file in file_names:
                attachment_path = os.path.join(attachment_folder, file)
                attachment_parsing(docs, attachment_path, file, metadata, text_splitter)

    elif filename.endswith(".txt"):
        loader = TextLoader(attachment_path)
        loaded = loader.load()
        for i, doc in enumerate(loaded):
            loaded[i].metadata.update(metadata)
        
        if isinstance(loaded[0], list):
            print(filename)
        
        docs.extend(loaded)

    else:
        pass


def export_data_win(items):
    make_folders()
    docs = list()
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    for _, message in enumerate(items):    
        body = unquote(message.Body)

        # metadata
        recipients_info = str()
        for idx in range(message.Recipients.Count): 
            recipient = message.Recipients.Item(idx+1)
            recipients_info += f"{recipient.Name} <{recipient.Address}>, "

        metadata = {
            "Subject": str(message.Subject),
            "From": f"{message.SenderName} <{message.SenderEmailAddress}>",
            "To": recipients_info,
            "Date": str(message.ReceivedTime.strftime("%a, %d %b %Y %H:%M:%S %z")),
            "Attchment": False
        }

        # docs for body
        docs.extend([Document(page_content=x, metadata = metadata) for x in text_splitter.split_text(body)])

        # docs for attachments
        for attachment in message.Attachments:
            attachment_path = os.path.join(attachment_folder, attachment.FileName)
            attachment.SaveAsFile(attachment_path)
            filename = attachment.FileName
            metadata['Attchment'] = filename
            attachment_parsing(docs, attachment_path, filename, metadata, text_splitter)

    return docs




if __name__ == "__main__":
    # export_data_pop()
    items = fetch_outlook()
    docs = export_data_win(items)

    # embeddings = OpenAIEmbeddings()

    # vectordb = Chroma.from_documents(documents=docs, 
    #                            embedding=embeddings,
    #                            collection_name="outlooks",
    #                            persist_directory="chroma_folder")
    
    # vectordb.persist()
    